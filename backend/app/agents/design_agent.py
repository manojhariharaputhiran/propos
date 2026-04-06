"""
Design Agent: generates the PPT presentation using python-pptx.
Supports custom uploaded templates with {{title}}, {{body_text}}, {{chart}} placeholders.
Falls back to a built-in professional default template.
"""
import io
import os
from typing import List, Dict, Any, Optional
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# Brand colours for default template
BRAND_PRIMARY = RGBColor(0x1E, 0x40, 0xAF)   # Deep blue
BRAND_ACCENT = RGBColor(0x06, 0xB6, 0xD4)    # Cyan
BRAND_TEXT = RGBColor(0x1F, 0x29, 0x37)      # Near-black
BRAND_LIGHT = RGBColor(0xF8, 0xFA, 0xFF)     # Off-white


def build_ppt(
    questions: List[Dict[str, Any]],
    title: str = "RFP Response",
    brand_voice: str = "professional",
    template_path: Optional[str] = None,
    critic_report: Optional[Dict[str, Any]] = None,
) -> bytes:
    """
    Build a PPT from questions/answers. Returns raw bytes.
    """
    if template_path and os.path.exists(template_path):
        return _build_from_template(questions, title, template_path, critic_report)
    return _build_default_ppt(questions, title, brand_voice, critic_report)


def _build_default_ppt(
    questions: List[Dict[str, Any]],
    title: str,
    brand_voice: str,
    critic_report: Optional[Dict[str, Any]],
) -> bytes:
    """Create a professional PPT from scratch."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, BRAND_PRIMARY)
    _add_text_box(
        slide,
        title,
        Inches(1), Inches(2.5), Inches(11.33), Inches(1.5),
        font_size=Pt(40), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF)
    )
    _add_text_box(
        slide,
        "Proposal Response Document",
        Inches(1), Inches(4.2), Inches(8), Inches(0.6),
        font_size=Pt(20), color=BRAND_ACCENT
    )

    # ── Slide 2: Scorecard (if critic report available) ───────────────────────
    if critic_report:
        slide = prs.slides.add_slide(blank_layout)
        _add_bg(slide, BRAND_LIGHT)
        _add_text_box(
            slide, "Proposal Scorecard",
            Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            font_size=Pt(28), bold=True, color=BRAND_PRIMARY
        )
        scores = [
            ("Win Probability", critic_report.get("win_probability", 0)),
            ("Completeness", critic_report.get("completeness", 0)),
            ("Persuasiveness", critic_report.get("persuasiveness", 0)),
            ("Compliance", critic_report.get("compliance", 0)),
        ]
        for i, (label, score) in enumerate(scores):
            x = Inches(0.5 + i * 3.2)
            _add_score_box(slide, label, score, x, Inches(1.5))

        suggestions = critic_report.get("suggestions", [])
        if suggestions:
            suggestions_text = "Key Improvements:\n" + "\n".join(f"• {s}" for s in suggestions[:3])
            _add_text_box(
                slide, suggestions_text,
                Inches(0.5), Inches(4.5), Inches(12), Inches(2.5),
                font_size=Pt(14), color=BRAND_TEXT
            )

    # ── Q&A Slides ────────────────────────────────────────────────────────────
    for i, q in enumerate(questions):
        slide = prs.slides.add_slide(blank_layout)
        _add_bg(slide, RGBColor(0xFF, 0xFF, 0xFF))

        # Slide number accent bar
        tf = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0), Inches(13.33), Inches(0.08)
        )
        tf.fill.solid()
        tf.fill.fore_color.rgb = BRAND_PRIMARY
        tf.line.fill.background()

        # Question
        _add_text_box(
            slide,
            f"Q{i + 1}: {q['text']}",
            Inches(0.5), Inches(0.3), Inches(12.33), Inches(1.2),
            font_size=Pt(16), bold=True, color=BRAND_PRIMARY
        )

        # Answer
        answer = q.get("final_answer") or q.get("draft") or "Pending SME input."
        _add_text_box(
            slide, answer,
            Inches(0.5), Inches(1.7), Inches(12.33), Inches(4.8),
            font_size=Pt(14), color=BRAND_TEXT
        )

        # Confidence badge
        conf = q.get("confidence")
        if conf is not None:
            conf_text = f"Confidence: {int(conf * 100)}%"
            conf_color = BRAND_ACCENT if conf >= 0.7 else RGBColor(0xF5, 0x9E, 0x0B)
            _add_text_box(
                slide, conf_text,
                Inches(10.5), Inches(6.8), Inches(2.5), Inches(0.4),
                font_size=Pt(10), color=conf_color
            )

    # ── Final Slide ───────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, BRAND_PRIMARY)
    _add_text_box(
        slide, "Thank You",
        Inches(3), Inches(3), Inches(7), Inches(1.5),
        font_size=Pt(48), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF)
    )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_from_template(
    questions: List[Dict[str, Any]],
    title: str,
    template_path: str,
    critic_report: Optional[Dict[str, Any]],
) -> bytes:
    """Populate a user-uploaded template with {{placeholder}} replacements."""
    prs = Presentation(template_path)

    def _replace_text(shape, replacements: Dict[str, str]):
        if not shape.has_text_frame:
            return
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for key, val in replacements.items():
                    if key in run.text:
                        run.text = run.text.replace(key, val)

    answers_block = "\n\n".join(
        f"Q{i+1}: {q['text']}\nA: {(q.get('final_answer') or q.get('draft') or 'TBD')}"
        for i, q in enumerate(questions)
    )

    replacements = {
        "{{title}}": title,
        "{{body_text}}": answers_block[:2000],
        "{{chart}}": "",  # Placeholder cleared; chart added separately
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            _replace_text(shape, replacements)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Helper drawing functions ──────────────────────────────────────────────────

def _add_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, text, left, top, width, height,
                  font_size=Pt(14), bold=False, color=None, align=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if align:
        p.alignment = align
    return txBox


def _add_score_box(slide, label: str, score: int, left, top):
    """Add a score metric box with label and value."""
    color = (
        RGBColor(0x16, 0xA3, 0x4A) if score >= 75 else
        RGBColor(0xF5, 0x9E, 0x0B) if score >= 50 else
        RGBColor(0xDC, 0x26, 0x26)
    )
    box = slide.shapes.add_shape(1, left, top, Inches(2.8), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.color.rgb = color
    box.line.width = Pt(2)

    txBox = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.2), Inches(2.6), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{score}%"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(1.4), Inches(2.6), Inches(0.7))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(12)
    run2.font.color.rgb = BRAND_TEXT
    p2.alignment = PP_ALIGN.CENTER
